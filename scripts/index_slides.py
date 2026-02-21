"""Index all slides from branding presentations."""
from pptx import Presentation
import os

bases = {
    'MOD1': r'C:\Users\vidal\OneDrive\Documentos\1 - ACADEMICO\8 - MEUS CURSOS\1 - INTRODUÇÃO AO BRANDING\AULAS\PRIMEIRO MÓDULO\AULA',
    'MOD2': r'C:\Users\vidal\OneDrive\Documentos\1 - ACADEMICO\8 - MEUS CURSOS\1 - INTRODUÇÃO AO BRANDING\AULAS\SEGUNDO MODULO\AULAS',
}

targets = [
    ('MOD1', '1. MUDANÇAS DE CENÁRIO DO MERCADO.pptx'),
    ('MOD1', '2. INTRODUÇÃO AO BRANDING.pptx'),
    ('MOD1', '5. ESTRATÉGIA DE PRODUTOS.pptx'),
    ('MOD1', '6. PÚBLICO IDEAL.pptx'),
    ('MOD2', '4. INTRODUÇÃO AO POSICIONAMENTO.pptx'),
]

out = open(r'scripts\slide_index.txt', 'w', encoding='utf-8')

for mod, fname in targets:
    path = os.path.join(bases[mod], fname)
    prs = Presentation(path)
    header = f'=== [{mod}] {fname} ({len(prs.slides)} slides) ==='
    print(header)
    out.write(header + '\n')
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            summary = ' | '.join(texts[:5])
            if len(summary) > 150:
                summary = summary[:150] + '...'
            line = f'  s{i:03d}: {summary}'
        else:
            line = f'  s{i:03d}: [sem texto]'
        out.write(line + '\n')
    out.write('\n')

out.close()
print('Done.')
