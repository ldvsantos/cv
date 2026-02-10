#!/usr/bin/env python3
"""Extract PPTX content from bioengenharia_de_solos and print for analysis."""

import sys
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

def clean_text(text):
    if not text:
        return ""
    replacements = {
        '\uf0ad': '→', '\uf0e8': '→', '\uf038': '•', '\uf0b7': '•',
        '\uf0a7': '§', '\uf0d8': '◆', '\uf076': '✓', '\uf0fc': '✓',
        '\uf0a8': '■', '\uf0b2': '►', '\uf06e': '●', '\uf0de': '▸',
        '\uf046': '✔', '\uf0d0': '–', '\uf02d': '-', '\uf0b0': '°',
        '\uf0e0': '→', '\uf0d2': '●', '\uf0a4': '→',
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text

def extract_and_print(filepath):
    try:
        prs = Presentation(str(filepath))
    except Exception as e:
        print(f"ERROR opening {filepath.name}: {e}")
        return
    
    print(f"\n{'='*80}")
    print(f"FILE: {filepath.name}")
    print(f"PATH: {filepath.parent.name}/{filepath.name}")
    print(f"SLIDES: {len(prs.slides)}")
    print(f"{'='*80}")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n--- SLIDE {slide_num} ---")
        has_content = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text_parts = []
                    is_title = "title" in shape.name.lower()
                    font_size = None
                    
                    for run in para.runs:
                        t = clean_text(run.text)
                        if not t:
                            continue
                        if run.font.size:
                            font_size = run.font.size
                        if run.font.bold:
                            t = f"**{t}**"
                        if run.font.italic:
                            t = f"*{t}*"
                        text_parts.append(t)
                    
                    full = "".join(text_parts).strip()
                    if full:
                        level = para.level or 0
                        big = font_size and font_size >= Emu(400000) if font_size else False
                        marker = "TITLE" if (is_title or big) else f"L{level}"
                        print(f"  [{marker}] {full}")
                        has_content = True
            
            if shape.has_table:
                table = shape.table
                print("  [TABLE]")
                for row in table.rows:
                    cells = [clean_text(cell.text.strip()) for cell in row.cells]
                    print(f"  | {' | '.join(cells)} |")
                has_content = True
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"  [IMAGE: {shape.name}]")
                has_content = True
        
        if not has_content:
            print("  (empty slide)")

BASE = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\bioengenharia_de_solos")

pptx_files = sorted(BASE.rglob("*.pptx"))
print(f"Found {len(pptx_files)} PPTX files:\n")
for f in pptx_files:
    print(f"  - {f.relative_to(BASE)}")

for f in pptx_files:
    extract_and_print(f)
