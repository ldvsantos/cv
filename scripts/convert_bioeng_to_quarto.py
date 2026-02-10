#!/usr/bin/env python3
"""Convert PPTX files from bioengenharia_de_solos to Quarto revealjs slides."""

import sys
import io
import re
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

BASE_REF = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\bioengenharia_de_solos")
OUTPUT_DIR = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\bioengenharia_de_solos\aulas")
ASSETS_SRC = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\ciencia_pi\SLIDES_QUARTO\Tema 01 - Gestao da Inovacao Tecnologica\assets")

# Mapping: (relative_path, folder_name, title, subtitle, order)
AULA_MAP = [
    ("Aulas 1 e 2/O solo.pptx",
     "aula01_o_solo",
     "O Solo",
     "Classificação e Propriedades dos Solos Brasileiros",
     1),
    ("Aulas 3 e 4/Aula 3. Solo X Erosão.pptx",
     "aula02_solo_erosao",
     "Solo e Erosão",
     "Processos Erosivos e Fatores Condicionantes",
     2),
    ("Aulas 3 e 4/Aula Prática 3 - Declividade - Equipamentos.pptx",
     "aula03_declividade_equipamentos",
     "Declividade – Equipamentos",
     "Aspectos Metodológicos e Instrumentos de Medição",
     3),
    ("Aulas 5, 6 , 7, 8 e 9/Controle de erosão hidrica.pptx",
     "aula04_controle_erosao_hidrica",
     "Controle de Erosão Hídrica",
     "Práticas Mecânicas e Terraceamento",
     4),
    ("Aulas 10  e 11 - Canal Escoadouro/Canal Escoadouro - Teoria.ptx.pptx",
     "aula05_canal_escoadouro_teoria",
     "Canal Escoadouro – Teoria",
     "Dimensionamento e Cálculo de Canais Escoadouros",
     5),
    ("Aulas 10  e 11 - Canal Escoadouro/Canal Escoadouro - Passo a passo.pptx",
     "aula06_canal_escoadouro_pratica",
     "Canal Escoadouro – Passo a Passo",
     "Exercício Prático de Cálculo e Construção",
     6),
]


def clean_text(text):
    if not text:
        return ""
    replacements = {
        '\uf0ad': '→', '\uf0e8': '→', '\uf038': '•', '\uf0b7': '•',
        '\uf0a7': '§', '\uf0d8': '◆', '\uf076': '✓', '\uf0fc': '✓',
        '\uf0a8': '■', '\uf0b2': '►', '\uf06e': '●', '\uf0de': '▸',
        '\uf046': '✔', '\uf0d0': '–', '\uf02d': '-', '\uf0b0': '°',
        '\uf0e0': '→', '\uf0d2': '●', '\uf0a4': '→',
        '\u00a0': ' ',  # non-breaking space
        '–': '–',
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    # Remove leftover unicode math symbols that didn't decode well
    text = re.sub(r'[\uf000-\uf0ff]', '', text)
    # Remove weird leftover chars from equation encoding
    text = re.sub(r'[Øæ¬Øæ░Øæ¿ØƒæØƒöØƒÄØÆöÔêÆØƒÅØÆÄØÆëØæáØæÜØÉ´ØæâM]+', '', text)
    return text


def extract_pptx_slides(filepath):
    """Extract slide content from PPTX file into structured data."""
    prs = Presentation(str(filepath))
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {"num": slide_num, "title": "", "content": [], "tables": [], "has_image": False, "image_count": 0}

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data["has_image"] = True
                slide_data["image_count"] += 1

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [clean_text(cell.text.strip()) for cell in row.cells]
                    rows.append(cells)
                if rows:
                    slide_data["tables"].append(rows)

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text_parts = []
                    is_title_shape = "title" in shape.name.lower() or "título" in shape.name.lower()
                    font_size = None
                    has_bold = False

                    for run in para.runs:
                        t = clean_text(run.text)
                        if not t.strip():
                            continue
                        if run.font.size:
                            font_size = run.font.size
                        if run.font.bold:
                            has_bold = True
                            t = f"**{t}**"
                        if run.font.italic:
                            t = f"*{t}*"
                        text_parts.append(t)

                    full_text = "".join(text_parts).strip()
                    # Clean up double asterisks
                    full_text = re.sub(r'\*\*\s*\*\*', '', full_text)
                    full_text = re.sub(r'\*\*\*\*', '**', full_text)
                    full_text = full_text.strip()
                    if not full_text:
                        continue

                    level = para.level or 0
                    is_big = font_size and font_size >= Emu(400000) if font_size else False

                    if is_title_shape or is_big:
                        if not slide_data["title"]:
                            clean_title = re.sub(r'\*+', '', full_text).strip()
                            slide_data["title"] = clean_title
                        else:
                            slide_data["content"].append({"text": full_text, "level": level, "is_bold": has_bold})
                    else:
                        slide_data["content"].append({"text": full_text, "level": level, "is_bold": has_bold})

        slides.append(slide_data)

    return slides


def slides_to_qmd(slides, title, subtitle, aula_name):
    """Convert extracted slides to Quarto revealjs QMD format."""
    lines = []

    # YAML header
    lines.append("---")
    lines.append(f'title: "{title}"')
    lines.append(f'subtitle: "{subtitle}<br>Bioengenharia de Solos"')
    lines.append('author: "Luiz Diego Vidal Santos"')
    lines.append('institute: "Universidade Estadual de Feira de Santana (UEFS)"')
    lines.append("format:")
    lines.append("  revealjs:")
    lines.append("    logo: ../../assets/logo-uefs.webp")
    lines.append(f'    footer: "UEFS | Bioengenharia de Solos | {title}"')
    lines.append("    slide-number: true")
    lines.append("    theme: [simple, assets/uefs.scss]")
    lines.append("    controls: true")
    lines.append("    width: 1280")
    lines.append("    height: 720")
    lines.append("    css: assets/custom.css")
    lines.append("    transition: slide")
    lines.append("    background-transition: fade")
    lines.append("    preview-links: auto")
    lines.append("execute:")
    lines.append("  echo: false")
    lines.append("---")
    lines.append("")

    prev_title = ""
    slide_count = 0

    for i, slide in enumerate(slides):
        # Skip slides that are pure image-only with no text
        has_text = bool(slide["title"] or slide["content"] or slide["tables"])
        if not has_text:
            continue

        slide_count += 1
        slide_title = slide["title"] if slide["title"] else ""

        # First meaningful slide -> section header
        if slide_count == 1 and slide_title:
            lines.append(f"# [{slide_title}]{{.section-header}}")
            lines.append("")
            if slide["content"] or slide["tables"]:
                lines.append("---")
                lines.append("")
                lines.append(f"## {slide_title} {{.smaller-text}}")
                lines.append("")
                _render_content(lines, slide)
                lines.append("")
            lines.append("---")
            lines.append("")
            prev_title = slide_title
            continue

        # Section break: title only, no content or tables
        if slide_title and not slide["content"] and not slide["tables"]:
            lines.append(f"# [{slide_title}]{{.section-header}}")
            lines.append("")
            lines.append("---")
            lines.append("")
            prev_title = slide_title
            continue

        # Regular slide with content
        if slide_title:
            display_title = slide_title
            if display_title == prev_title:
                display_title = f"{slide_title} (cont.)"
            lines.append(f"## {display_title} {{.smaller-text}}")
            prev_title = slide_title
        else:
            lines.append("## {.smaller-text}")

        lines.append("")

        # Note: indicate images if present
        if slide["has_image"] and slide["image_count"] > 0:
            lines.append("::: {.callout-note}")
            lines.append(f"Slide original contém {slide['image_count']} imagem(ns) ilustrativa(s).")
            lines.append(":::")
            lines.append("")

        _render_content(lines, slide)
        lines.append("")
        lines.append("---")
        lines.append("")

    # Final slide
    lines.append('# {background-color="#034EA2"}')
    lines.append("")
    lines.append('[Obrigado!]{.obrigado-fit-text style="color: white;"}')
    lines.append("")
    lines.append('::: {style="text-align: center; color: white; font-size: 0.7em;"}')
    lines.append("Luiz Diego Vidal Santos")
    lines.append("")
    lines.append("Universidade Estadual de Feira de Santana (UEFS)")
    lines.append(":::")
    lines.append("")

    return "\n".join(lines)


def _render_content(lines, slide):
    """Render slide content items."""
    for item in slide["content"]:
        text = item["text"]
        level = item["level"]

        if level == 0:
            if text.startswith(("•", "→", "►", "▸", "●", "■", "-")):
                text = text.lstrip("•→►▸●■- ").strip()
                lines.append(f"- {text}")
            else:
                lines.append(text)
                lines.append("")
        elif level >= 1:
            indent = "  " * level
            if text.startswith(("•", "→", "►", "▸", "●", "■", "-")):
                text = text.lstrip("•→►▸●■- ").strip()
            lines.append(f"{indent}- {text}")

    # Render tables
    for table in slide["tables"]:
        lines.append("")
        if len(table) >= 1:
            header = table[0]
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in table[1:]:
                while len(row) < len(header):
                    row.append("")
                lines.append("| " + " | ".join(row[:len(header)]) + " |")
        lines.append("")


def main():
    # Create output directory
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Create assets directory and copy assets
    assets_dir = OUTPUT_DIR / "assets"
    assets_dir.mkdir(exist_ok=True)

    # Copy scss and css from ciencia_pi template
    for f in ["ufs.scss", "custom.css", "fundo.png", "contracapa.jpg"]:
        src = ASSETS_SRC / f
        dst_name = "uefs.scss" if f == "ufs.scss" else f
        if src.exists():
            shutil.copy2(src, assets_dir / dst_name)
            print(f"  Copied {f} -> {dst_name}")

    # Process each PPTX
    success_count = 0
    error_count = 0

    for rel_path, folder_name, title, subtitle, order in AULA_MAP:
        pptx_path = BASE_REF / rel_path

        if not pptx_path.exists():
            print(f"  SKIP (not found): {rel_path}")
            error_count += 1
            continue

        try:
            slides = extract_pptx_slides(pptx_path)
            qmd_content = slides_to_qmd(slides, title, subtitle, folder_name)
            qmd_file = OUTPUT_DIR / f"{folder_name}.qmd"
            qmd_file.write_text(qmd_content, encoding='utf-8')
            slide_count = sum(1 for s in slides if s["title"] or s["content"] or s["tables"])
            print(f"  OK [{order:02d}] {folder_name}.qmd ({slide_count} slides com conteúdo de {len(slides)} total)")
            success_count += 1
        except Exception as e:
            print(f"  ERROR: {rel_path} -> {e}")
            import traceback
            traceback.print_exc()
            error_count += 1

    print(f"\nDone: {success_count} converted, {error_count} errors")
    print(f"Output: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
