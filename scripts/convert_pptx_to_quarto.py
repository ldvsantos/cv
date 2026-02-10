#!/usr/bin/env python3
"""Convert all PPTX files from analise_dados_ambientais to Quarto revealjs slides."""

import sys
import io
import re
import os
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Force UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

BASE_REF = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\analise_dados_ambientais\referencias")
OUTPUT_DIR = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\analise_dados_ambientais\aulas")
ASSETS_SRC = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\ciencia_pi\SLIDES_QUARTO\Tema 01 - Gestao da Inovacao Tecnologica\assets")

# Mapping of PPTX file -> (folder_name, title, subtitle, order)
AULA_MAP = [
    # Pasta 1 - Introdução
    ("1 - INTRODUÇÃO ESTATÍSTICA DESCRITIVA/1 - INTRODUÇÃO ((ETIMOLOGIA)).pptx",
     "aula01_introducao_etimologia", "Introdução e Etimologia Estatística", "Conceitos Fundamentais de Estatística", 1),
    ("1 - INTRODUÇÃO ESTATÍSTICA DESCRITIVA/1 - INTRODUÇÃO A ESTATÍSTICA DESCRITIVA.pptx",
     "aula02_estatistica_descritiva", "Estatística Descritiva", "Medidas de Tendência Central e Dispersão", 2),
    # Pasta 2 - Inferencial
    ("2 - ESTATÍSTICA INFERENCIAL E DN/2 - ESTATISTICA INFERENCIAL E DISTRIBUIÇÃO NORMAL.pptx",
     "aula03_inferencial_dist_normal", "Estatística Inferencial e Distribuição Normal", "Fundamentos da Inferência Estatística", 3),
    ("2 - ESTATÍSTICA INFERENCIAL E DN/3 - ESCOLHA DO TESTE.pptx",
     "aula04_escolha_do_teste", "A Escolha do Teste Estatístico", "Como Selecionar a Análise Adequada", 4),
    ("2 - ESTATÍSTICA INFERENCIAL E DN/4 - TAXONOMIA DE BLOOM.pptx",
     "aula05_taxonomia_bloom", "Taxonomia de Bloom", "Níveis de Aprendizagem e Pesquisa", 5),
    # Pasta 3 - Correlação
    ("3 - ANÁLISE DE CORRELAÇÃO/3 - ANÁLISE DE CORRELAÇÃO.pptx",
     "aula06_correlacao", "Análise de Correlação", "Correlação de Pearson, Spearman e Kendall", 6),
    # Pasta 4 - Teste T
    ("4 - TESTE T/4 - TESTES DE HIPÓTESES - TESTE T.pptx",
     "aula07_teste_t", "Testes de Hipóteses - Teste T", "Teste T para Amostras Independentes e Pareadas", 7),
    # Pasta 5 - ANOVA
    ("5 - ANOVA/5.1 - TESTES DE HIPÓTESES - ANOVA UMA VIA.pptx",
     "aula08_anova_uma_via", "ANOVA de Uma Via", "Análise de Variância One-Way", 8),
    ("5 - ANOVA/5.2 - BOOTSTRAPPING.pptx",
     "aula09_bootstrapping", "Bootstrapping", "Técnicas de Reamostragem", 9),
    ("5 - ANOVA/5.3 - TESTES DE HIPÓTESES - ANOVA FATORIAL.pptx",
     "aula10_anova_fatorial", "ANOVA Fatorial", "Análise de Variância com Múltiplos Fatores", 10),
    ("5 - ANOVA/5 - TESTES DE HIPÓTESES - ANOVA COMPLETO.pptx",
     "aula11_anova_completo", "ANOVA - Visão Completa", "Análise de Variância Aprofundada", 11),
    ("5 - ANOVA/5 - TESTES DE HIPÓTESES - ANOVA MR.pptx",
     "aula12_anova_medidas_repetidas", "ANOVA de Medidas Repetidas", "Análise para Delineamentos com Medidas Repetidas", 12),
    # Pasta 6 - ANCOVA
    ("6 - ANCOVA/6 - TESTES DE HIPÓTESES - ANCOVA.pptx",
     "aula13_ancova", "ANCOVA", "Análise de Covariância", 13),
    # Pasta 8 - Metanálise
    ("8 - METANALISE/ARTIGO METANALISE.pptx",
     "aula14_metanalise", "Metanálise", "Fundamentos de Metanálise", 14),
    # Pasta 9 - Regressão
    ("9 - REGRESSÃO/REGRESSÃO.pptx",
     "aula15_regressao", "Regressão", "Modelos de Regressão Linear e Múltipla", 15),
    ("9 - REGRESSÃO/TESTES DE QUI-QUADRADO.pptx",
     "aula16_qui_quadrado", "Testes de Qui-Quadrado", "Teste de Independência e Aderência", 16),
    # Pasta 10 - GLM e GEE
    ("10 - GEE GLM/Slides GLM e GEE.pptx",
     "aula17_glm_gee", "GLM e GEE", "Modelos Lineares Generalizados e Equações de Estimação Generalizadas", 17),
    # Pasta 11 - Instrumentos
    ("11 - CONSTRUÇÃO E ADAPTAÇÃO INSTRU/Cons de Medidas.pptx",
     "aula18_construcao_instrumentos", "Construção de Instrumentos de Medida", "Desenvolvimento de Escalas e Questionários", 18),
    ("11 - CONSTRUÇÃO E ADAPTAÇÃO INSTRU/Adpt de Medidas.pptx",
     "aula19_adaptacao_instrumentos", "Adaptação de Instrumentos de Medida", "Tradução e Adaptação Transcultural", 19),
    ("11 - CONSTRUÇÃO E ADAPTAÇÃO INSTRU/Evidencias validade.pptx",
     "aula20_evidencias_validade", "Evidências de Validade", "Validade de Conteúdo, Construto e Critério", 20),
    ("11 - CONSTRUÇÃO E ADAPTAÇÃO INSTRU/Validade instrumento.pptx",
     "aula21_validade_instrumento", "Validade de Instrumentos", "Propriedades Psicométricas", 21),
    # Pasta 12 - AFE e AFC
    ("12 - AFE e AFC/AFE.pptx",
     "aula22_analise_fatorial", "Análise Fatorial Exploratória", "AFE - Fundamentos e Aplicações", 22),
    # Pasta 13 - TRI
    ("13 - TRI/introdução a Teoria de Resposta ao Item.pptx",
     "aula23_tri_introducao", "Introdução à TRI", "Teoria de Resposta ao Item", 23),
    ("13 - TRI/TRI RASCH.pptx",
     "aula24_tri_rasch", "Modelo de Rasch", "TRI - Modelo de Rasch", 24),
    # Raiz - arquivos avulsos
    ("ANÁLISE BI E MULTIVARIADA (COMPLETO).pptx",
     "aula25_analise_bi_multivariada", "Análise Bi e Multivariada", "Técnicas Multivariadas Completas", 25),
    ("BOOT.pptx",
     "aula26_bootstrap_avancado", "Bootstrap Avançado", "Técnicas Avançadas de Reamostragem", 26),
    ("MODELOS DE ASSOCIAÇÃO E DEPENDÊNCIA.pptx",
     "aula27_associacao_dependencia", "Modelos de Associação e Dependência", "Relações entre Variáveis", 27),
    ("NÃO PARAMÉTRICOS.pptx",
     "aula28_nao_parametricos", "Testes Não Paramétricos", "Alternativas aos Testes Paramétricos", 28),
]


def clean_text(text):
    """Clean extracted text, fixing common encoding issues from PPTX."""
    if not text:
        return ""
    # Replace common wingdings/symbol font chars
    replacements = {
        '\uf0ad': '→',
        '\uf0e8': '→',
        '\uf038': '•',
        '\uf0b7': '•',
        '\uf0a7': '§',
        '\uf0d8': '◆',
        '\uf076': '✓',
        '\uf0fc': '✓',
        '\uf0a8': '■',
        '\uf0b2': '►',
        '\uf06e': '●',
        '\uf0de': '▸',
        '\uf046': '✔',
        '\uf0d0': '–',
        '\uf02d': '-',
        '\uf0b0': '°',
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


def extract_pptx_slides(filepath):
    """Extract slide content from PPTX file into structured data."""
    prs = Presentation(str(filepath))
    slides = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {"num": slide_num, "title": "", "content": [], "tables": [], "has_image": False}
        
        for shape in slide.shapes:
            # Check for images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data["has_image"] = True
            
            # Check for tables
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [clean_text(cell.text.strip()) for cell in row.cells]
                    rows.append(cells)
                if rows:
                    slide_data["tables"].append(rows)
            
            # Text content
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text_parts = []
                    is_title_shape = "title" in shape.name.lower()
                    font_size = None
                    has_bold = False
                    
                    for run in para.runs:
                        t = clean_text(run.text)
                        if not t:
                            continue
                        if run.font.size:
                            font_size = run.font.size
                        if run.font.bold:
                            has_bold = True
                            t = f"**{t}**"
                        if run.font.italic:
                            t = f"*{t}*"
                        text_parts.append(t)
                    
                    full_text = "".join(text_parts).strip()
                    if not full_text:
                        continue
                    
                    level = para.level or 0
                    is_big = font_size and font_size >= Emu(500000) if font_size else False
                    
                    # Detect title
                    if is_title_shape or is_big:
                        if not slide_data["title"]:
                            # Clean markdown from title
                            clean_title = re.sub(r'\*+', '', full_text).strip()
                            slide_data["title"] = clean_title
                        else:
                            slide_data["content"].append({"text": full_text, "level": level, "is_bold": has_bold})
                    else:
                        slide_data["content"].append({"text": full_text, "level": level, "is_bold": has_bold})
        
        slides.append(slide_data)
    
    return slides


def slides_to_qmd(slides, title, subtitle, aula_name):
    """Convert extracted slides to Quarto revealjs QMD format."""
    lines = []
    
    # YAML header
    lines.append("---")
    lines.append(f'title: "{title}"')
    lines.append(f'subtitle: "{subtitle}<br>Análise de Dados Ambientais"')
    lines.append('author: "Luiz Diego Vidal Santos"')
    lines.append('institute: "Universidade Estadual de Feira de Santana (UEFS)"')
    lines.append("format:")
    lines.append("  revealjs:")
    lines.append("    logo: ../../../assets/logo-uefs.webp")
    lines.append(f'    footer: "UEFS | Análise de Dados Ambientais | {title}"')
    lines.append("    slide-number: true")
    lines.append("    theme: [simple, assets/uefs.scss]")
    lines.append("    controls: true")
    lines.append("    width: 1280")
    lines.append("    height: 720")
    lines.append("    css: assets/custom.css")
    lines.append("    transition: slide")
    lines.append("    background-transition: fade")
    lines.append("    preview-links: auto")
    lines.append("execute:")
    lines.append("  echo: false")
    lines.append("---")
    lines.append("")
    
    prev_title = ""
    
    for i, slide in enumerate(slides):
        if not slide["title"] and not slide["content"] and not slide["tables"]:
            continue
        
        slide_title = slide["title"] if slide["title"] else ""
        
        # Section header for level 1 headings (first slide or major sections)
        if i == 0 and slide_title:
            lines.append(f"# [{slide_title}]{{.section-header}}")
            lines.append("")
            lines.append("---")
            lines.append("")
            # Also output content if present
            if slide["content"]:
                lines.append(f"## {slide_title} {{.smaller-text}}")
                lines.append("")
                _render_content(lines, slide)
                lines.append("")
                lines.append("---")
                lines.append("")
            prev_title = slide_title
            continue
        
        # Check if this is a section break (title only, no content)
        if slide_title and not slide["content"] and not slide["tables"]:
            lines.append(f"# [{slide_title}]{{.section-header}}")
            lines.append("")
            lines.append("---")
            lines.append("")
            prev_title = slide_title
            continue
        
        # Regular slide
        if slide_title:
            # If same title as previous, add subtitle variant
            display_title = slide_title
            if display_title == prev_title:
                display_title = f"{slide_title} (cont.)"
            lines.append(f"## {display_title} {{.smaller-text}}")
            prev_title = slide_title
        else:
            lines.append("## {.smaller-text}")
        
        lines.append("")
        _render_content(lines, slide)
        lines.append("")
        lines.append("---")
        lines.append("")
    
    # Final slide
    lines.append("# {background-color=\"#034EA2\"}")
    lines.append("")
    lines.append('[Obrigado!]{.obrigado-fit-text style="color: white;"}')
    lines.append("")
    lines.append("::: {style=\"text-align: center; color: white; font-size: 0.7em;\"}")
    lines.append("Luiz Diego Vidal Santos")
    lines.append("")
    lines.append("Universidade Estadual de Feira de Santana (UEFS)")
    lines.append(":::")
    lines.append("")
    
    return "\n".join(lines)


def _render_content(lines, slide):
    """Render slide content items."""
    for item in slide["content"]:
        text = item["text"]
        level = item["level"]
        
        if level == 0:
            # Check if it looks like a bullet point
            if text.startswith(("•", "→", "►", "▸", "●", "■", "-")):
                text = text.lstrip("•→►▸●■- ").strip()
                lines.append(f"- {text}")
            else:
                lines.append(text)
                lines.append("")
        elif level >= 1:
            indent = "  " * level
            if text.startswith(("•", "→", "►", "▸", "●", "■", "-")):
                text = text.lstrip("•→►▸●■- ").strip()
            lines.append(f"{indent}- {text}")
    
    # Render tables
    for table in slide["tables"]:
        lines.append("")
        if len(table) >= 1:
            # Header row
            header = table[0]
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            # Data rows
            for row in table[1:]:
                # Pad row to match header length
                while len(row) < len(header):
                    row.append("")
                lines.append("| " + " | ".join(row[:len(header)]) + " |")
        lines.append("")


def main():
    # Create output directory
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    # Create assets directory and copy assets
    assets_dir = OUTPUT_DIR / "assets"
    assets_dir.mkdir(exist_ok=True)
    
    # Copy scss and css from ciencia_pi template
    for f in ["uefs.scss", "custom.css"]:
        src = ASSETS_SRC / f
        if src.exists():
            shutil.copy2(src, assets_dir / f)
            print(f"  Copied {f} to assets/")
    
    # Copy background images if present
    for f in ["fundo.png", "contracapa.jpg"]:
        src = ASSETS_SRC / f
        if src.exists():
            shutil.copy2(src, assets_dir / f)
            print(f"  Copied {f} to assets/")
    
    # Process each PPTX
    success_count = 0
    error_count = 0
    
    for rel_path, folder_name, title, subtitle, order in AULA_MAP:
        pptx_path = BASE_REF / rel_path
        
        if not pptx_path.exists():
            print(f"  SKIP (not found): {rel_path}")
            error_count += 1
            continue
        
        try:
            # Extract slides
            slides = extract_pptx_slides(pptx_path)
            
            # Generate QMD
            qmd_content = slides_to_qmd(slides, title, subtitle, folder_name)
            
            # Write to output
            qmd_file = OUTPUT_DIR / f"{folder_name}.qmd"
            qmd_file.write_text(qmd_content, encoding='utf-8')
            
            print(f"  OK [{order:02d}] {folder_name}.qmd ({len(slides)} slides)")
            success_count += 1
            
        except Exception as e:
            print(f"  ERROR: {rel_path} -> {e}")
            error_count += 1
    
    print(f"\nDone: {success_count} converted, {error_count} errors")
    print(f"Output: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
