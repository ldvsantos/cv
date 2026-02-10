#!/usr/bin/env python3
"""Extract PPTX content and print structured text for Quarto conversion."""

import sys
import json
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Force UTF-8 for stdout
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

def extract_pptx(filepath):
    """Extract all slides content from a PPTX file."""
    prs = Presentation(filepath)
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "number": slide_num,
            "shapes": [],
            "layout": slide.slide_layout.name if slide.slide_layout else "unknown"
        }
        
        for shape in slide.shapes:
            shape_data = {
                "type": shape.shape_type,
                "name": shape.name,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            
            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    para_text = ""
                    is_bold = False
                    is_italic = False
                    font_size = None
                    level = para.level if para.level else 0
                    
                    for run in para.runs:
                        text = run.text
                        if run.font.bold:
                            text = f"**{text}**"
                            is_bold = True
                        if run.font.italic:
                            text = f"*{text}*"
                            is_italic = True
                        if run.font.size:
                            font_size = run.font.size
                        para_text += text
                    
                    if para_text.strip():
                        paragraphs.append({
                            "text": para_text.strip(),
                            "level": level,
                            "font_size": str(font_size) if font_size else None,
                            "alignment": str(para.alignment) if para.alignment else None
                        })
                
                if paragraphs:
                    shape_data["paragraphs"] = paragraphs
                    slide_info["shapes"].append(shape_data)
            
            elif shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(cell.text.strip())
                    rows.append(cells)
                shape_data["table"] = rows
                slide_info["shapes"].append(shape_data)
        
        if slide_info["shapes"]:
            slides_data.append(slide_info)
    
    return slides_data


def slides_to_text(slides_data, filename):
    """Convert slides data to readable text format."""
    lines = []
    lines.append(f"=== FILE: {filename} ===")
    lines.append(f"Total slides: {len(slides_data)}")
    lines.append("")
    
    for slide in slides_data:
        lines.append(f"--- SLIDE {slide['number']} (layout: {slide['layout']}) ---")
        for shape in slide["shapes"]:
            if "paragraphs" in shape:
                for p in shape["paragraphs"]:
                    prefix = "  " * p["level"]
                    size_info = f" [size:{p['font_size']}]" if p['font_size'] else ""
                    lines.append(f"{prefix}[L{p['level']}]{size_info} {p['text']}")
            if "table" in shape:
                lines.append("[TABLE]")
                for row in shape["table"]:
                    lines.append("| " + " | ".join(row) + " |")
        lines.append("")
    
    return "\n".join(lines)


if __name__ == "__main__":
    base = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\analise_dados_ambientais\referencias")
    
    # Collect all pptx files
    pptx_files = sorted(base.rglob("*.pptx"))
    
    for pptx_file in pptx_files:
        try:
            slides = extract_pptx(str(pptx_file))
            text = slides_to_text(slides, pptx_file.name)
            print(text)
            print("\n" + "=" * 80 + "\n")
        except Exception as e:
            print(f"ERROR processing {pptx_file.name}: {e}")
            print()
