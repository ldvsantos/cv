#!/usr/bin/env python3
"""Convert PPTX files of the Geotêxteis course to Quarto revealjs slides."""

import sys
import io
import re
import os
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Force UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# Base path for PPTX source files for Geotêxteis course
BASE_REF = Path(r"C:\Users\vidal\OneDrive\Documentos\1 - ACADEMICO\2 - UEFS\1 - DISCIPLINAS\4 - TOPICOS ESPECIAIS GEOTÊXTIL\1 - SLIDES")
# Output directory under site aulas/geotexteis
OUTPUT_DIR = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\geotexteis")
# Assets from standard template (reuse from ciencia_pi)
ASSETS_SRC = Path(r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\aulas\ciencia_pi\SLIDES_QUARTO\Tema 01 - Gestao da Inovacao Tecnologica\assets")

# Map PPTX filename -> (slug, title, subtitle, order)
AULA_MAP = [
    ("Apresentação POSDOC.pptx",
     "aula_geotextil_posdoc",
     "Seminário POSDOC: Tópicos Especiais em Geotêxteis",
     "Visão geral e aplicações",
     1),
    ("Aula Tipologia (aula 1).pptx",
     "aula_geotextil_tipologia",
     "Tipologia de Geotêxteis",
     "Classificação e propriedades",
     2),
    ("Aula Geotexteis resumida.pptx",
     "aula_geotextil_resumida",
     "Geotêxteis: Resumo abrangente",
     "Tipos e usos práticos",
     3),
    ("Aula Projeto (aula 2).pptx",
     "aula_geotextil_projeto2",
     "Projeto Prático I",
     "Estudo de caso em geotêxteis",
     4),
    ("Aula Projeto (aula 3).pptx",
     "aula_geotextil_projeto3",
     "Projeto Prático II",
     "Implementação e avaliação",
     5),
]


def clean_text(text):
    if not text:
        return ""
    replacements = {
        '\uf0ad': '→', '\uf0e8': '→', '\uf038': '•', '\uf0b7': '•',
        '\uf0a7': '§', '\uf0d8': '◆', '\uf076': '✓', '\uf0fc': '✓',
        '\uf0a8': '■', '\uf0b2': '►', '\uf06e': '●', '\uf0de': '▸',
        '\uf046': '✔', '\uf0d0': '–', '\uf02d': '-', '\uf0b0': '°',
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


def extract_pptx_slides(filepath):
    prs = Presentation(str(filepath))
    slides = []
    for slide in prs.slides:
        slide_data = {"num": None, "title": "", "content": [], "tables": [], "has_image": False}
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data["has_image"] = True
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([clean_text(c.text.strip()) for c in row.cells])
                if rows:
                    slide_data["tables"].append(rows)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text_parts = []
                    font_size = None
                    has_bold = False
                    for run in para.runs:
                        t = clean_text(run.text)
                        if not t:
                            continue
                        if run.font.bold:
                            has_bold = True
                            t = f"**{t}**"
                        text_parts.append(t)
                    full_text = "".join(text_parts).strip()
                    if not full_text:
                        continue
                    is_title = has_bold and not slide_data["title"]
                    if is_title:
                        slide_data["title"] = re.sub(r'\*+', '', full_text)
                    else:
                        slide_data["content"].append({"text": full_text, "level": para.level or 0})
        slides.append(slide_data)
    return slides


def slides_to_qmd(slides, title, subtitle, aula_name):
    lines = []
    lines.append("---")
    lines.append(f'title: "{title}"')
    lines.append(f'subtitle: "{subtitle}<br>Curso de Geotêxteis"')
    lines.append('author: "Luiz Diego Vidal Santos"')
    lines.append('institute: "UEFS - Universidade Estadual de Feira de Santana"')
    lines.append("format:")
    lines.append("  revealjs:")
    lines.append("    logo: ../../../assets/logo-uefs.webp")
    lines.append(f'    footer: "UEFS | Geotêxteis | {title}"')
    lines.append("    slide-number: true")
    lines.append("    theme: [simple, assets/uefs.scss]")
    lines.append("    controls: true")
    lines.append("    width: 1280")
    lines.append("    height: 720")
    lines.append("    css: assets/custom.css")
    lines.append("    transition: slide")
    lines.append("execute:")
    lines.append("  echo: false")
    lines.append("---\n")
    prev_title = ""
    for slide in slides:
        title_slide = slide["title"]
        if title_slide and title_slide != prev_title:
            lines.append(f"## {title_slide}{{.section-header}}\n")
            prev_title = title_slide
        for item in slide["content"]:
            text = item["text"]
            lines.append(f"- {text}")
        if slide["tables"]:
            for table in slide["tables"]:
                if table:
                    header = table[0]
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join(["---"]*len(header)) + " |")
                    for row in table[1:]:
                        lines.append("| " + " | ".join(row) + " |")
        lines.append("---\n")
    lines.append("# {background-color=\"#034EA2\"}\n")
    lines.append("[Obrigado!]{.obrigado-fit-text style=\"color: white;\"}\n")
    lines.append("::: {style=\"text-align: center; color: white; font-size: 0.7em;\"}\nUEFS - Geotêxteis\n:::\n")
    return "\n".join(lines)


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    assets_dir = OUTPUT_DIR / "assets"
    assets_dir.mkdir(exist_ok=True)
    for f in ["uefs.scss", "custom.css"]:
        src = ASSETS_SRC / f
        if src.exists():
            shutil.copy2(src, assets_dir / f)
    success = 0
    errors = 0
    for rel, slug, title, subtitle, order in AULA_MAP:
        pptx_path = BASE_REF / rel
        if not pptx_path.exists():
            print(f"SKIP not found: {rel}")
            errors += 1
            continue
        try:
            slides = extract_pptx_slides(pptx_path)
            qmd = slides_to_qmd(slides, title, subtitle, slug)
            outfile = OUTPUT_DIR / f"{slug}.qmd"
            outfile.write_text(qmd, encoding='utf-8')
            print(f"OK {slug}.qmd ({len(slides)} slides)")
            success += 1
        except Exception as e:
            print(f"ERROR {rel}: {e}")
            errors += 1
    print(f"Converted: {success}, Errors: {errors}")

if __name__ == "__main__":
    main()
