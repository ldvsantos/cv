"""
Extrai imagens embutidas dos arquivos PPTX de branding.
Salva cada imagem encontrada em books/branding-agro/img/slides/
"""
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

slides_dir = r"C:\Users\vidal\OneDrive\Documentos\1 - ACADEMICO\8 - MEUS CURSOS\1 - INTRODUÇÃO AO BRANDING\AULAS\PRIMEIRO MÓDULO\AULA"
out_dir = r"C:\Users\vidal\OneDrive\Documentos\13 - CLONEGIT\meu_site\books\branding-agro\img\slides"
os.makedirs(out_dir, exist_ok=True)

pptx_files = [
    "1. MUDANÇAS DE CENÁRIO DO MERCADO.pptx",
    "2. INTRODUÇÃO AO BRANDING.pptx",
    "3. INTRODUÇÃO A ESTRATÉGIA E DIFERENCIAÇÃO.pptx",
    "4. IDEIAS DIFERENCIADORAS.pptx",
    "5. ESTRATÉGIA DE PRODUTOS.pptx",
    "6. PÚBLICO IDEAL.pptx",
]

for pptx_name in pptx_files:
    pptx_path = os.path.join(slides_dir, pptx_name)
    if not os.path.exists(pptx_path):
        print(f"SKIP (not found): {pptx_name}")
        continue
    
    prefix = pptx_name.split(".")[0].strip()
    print(f"\n=== {pptx_name} ===")
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ERROR opening: {e}")
        continue
    
    img_count = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                img_count += 1
                fname = f"slide{prefix}_s{slide_idx:02d}_img{img_count:02d}.{ext}"
                fpath = os.path.join(out_dir, fname)
                with open(fpath, "wb") as f:
                    f.write(image.blob)
                print(f"  Slide {slide_idx}: {fname} ({len(image.blob)//1024}KB)")
            
            # Check grouped shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = child.image
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        img_count += 1
                        fname = f"slide{prefix}_s{slide_idx:02d}_img{img_count:02d}.{ext}"
                        fpath = os.path.join(out_dir, fname)
                        with open(fpath, "wb") as f:
                            f.write(image.blob)
                        print(f"  Slide {slide_idx} (group): {fname} ({len(image.blob)//1024}KB)")
    
    print(f"  Total: {img_count} imagens extraídas")

# Also check segundo modulo
slides_dir2 = slides_dir.replace("PRIMEIRO MÓDULO\\AULA", "SEGUNDO MODULO\\AULAS")
pptx_files2 = ["4. INTRODUÇÃO AO POSICIONAMENTO.pptx"]
for pptx_name in pptx_files2:
    pptx_path = os.path.join(slides_dir2, pptx_name)
    if not os.path.exists(pptx_path):
        print(f"SKIP (not found): {pptx_name}")
        continue
    prefix = "7"
    print(f"\n=== {pptx_name} ===")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ERROR opening: {e}")
        continue
    img_count = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                img_count += 1
                fname = f"slide{prefix}_s{slide_idx:02d}_img{img_count:02d}.{ext}"
                fpath = os.path.join(out_dir, fname)
                with open(fpath, "wb") as f:
                    f.write(image.blob)
                print(f"  Slide {slide_idx}: {fname} ({len(image.blob)//1024}KB)")
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = child.image
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        img_count += 1
                        fname = f"slide{prefix}_s{slide_idx:02d}_img{img_count:02d}.{ext}"
                        fpath = os.path.join(out_dir, fname)
                        with open(fpath, "wb") as f:
                            f.write(image.blob)
                        print(f"  Slide {slide_idx} (group): {fname} ({len(image.blob)//1024}KB)")
    print(f"  Total: {img_count} imagens extraídas")

print("\n=== DONE ===")
print(f"Pasta: {out_dir}")
print(f"Total de arquivos: {len(os.listdir(out_dir))}")
